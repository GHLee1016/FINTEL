"""
기존 fintel_results.pptx 에 'ML 학습 방식 + 하이퍼파라미터 튜닝' 슬라이드를
마지막에 append. 기존 슬라이드는 손대지 않음.
"""
from __future__ import annotations
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

PROJECT_ROOT = Path(__file__).resolve().parent.parent
PPTX_PATH = PROJECT_ROOT / 'report' / 'fintel_results.pptx'

# 기존 PPT의 디자인 토큰과 동일하게 유지
ACCENT       = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT_LIGHT = RGBColor(0xE8, 0xEE, 0xF4)
HEADER_BG    = RGBColor(0x2A, 0x3D, 0x52)
HEADER_FG    = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT_BG   = RGBColor(0xF7, 0xF8, 0xFA)
TXT          = RGBColor(0x1F, 0x2A, 0x3A)
TXT_MUTED    = RGBColor(0x6B, 0x72, 0x80)
DIVIDER_GRAY = RGBColor(0xE5, 0xE7, 0xEB)


# ── 헬퍼 ───────────────────────────────
def _set_text(shape, text, size=14, bold=False, color=TXT,
               align=PP_ALIGN.LEFT, italic=False, font='Calibri'):
    tf = shape.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font


def _add_text_lines(shape, lines, size=12, bold=False, color=TXT,
                     align=PP_ALIGN.LEFT, spacing=4, font='Calibri'):
    tf = shape.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(spacing)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font


def _accent_bar(slide, top, height, width_inch=0.08):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), top,
                                  Inches(width_inch), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar


def _hr_line(slide, y, sw, left=Inches(0.5), right_offset=Inches(0.5)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, y,
                                   sw - left - right_offset, Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = DIVIDER_GRAY
    line.line.fill.background()
    return line


def _title_block(slide, title, subtitle, sw):
    _accent_bar(slide, Inches(0.45), Inches(0.65))
    tx = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), sw - Inches(1.5), Inches(0.55))
    _set_text(tx, title, size=22, bold=True, color=TXT)
    sub = slide.shapes.add_textbox(Inches(0.75), Inches(0.92), sw - Inches(1.5), Inches(0.35))
    _set_text(sub, subtitle, size=12, color=TXT_MUTED, italic=True)
    _hr_line(slide, Inches(1.35), sw)


# ── 로드 + slide append ────────────────────────────
prs = Presentation(PPTX_PATH)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

slide = prs.slides.add_slide(BLANK)
_title_block(slide,
    'ML 학습 방식 & 하이퍼파라미터 튜닝',
    '5 ML 모델 (Ridge / EN / Huber / LightGBM / XGBoost) 의 학습·평가 절차',
    SW)

# ── 좌측: 학습 절차 ────────────────────
left_left = Inches(0.75)
left_top  = Inches(1.65)
left_w    = Inches(5.9)

label_l = slide.shapes.add_textbox(left_left, left_top, left_w, Inches(0.35))
_set_text(label_l, '1.  ML 학습 절차', size=13, bold=True, color=ACCENT)

# 박스
box_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                left_left, left_top + Inches(0.45),
                                left_w, Inches(4.3))
box_l.fill.solid()
box_l.fill.fore_color.rgb = ACCENT_LIGHT
box_l.line.color.rgb = ACCENT
box_l.line.width = Pt(1.0)

body_l = slide.shapes.add_textbox(left_left + Inches(0.25),
                                   left_top + Inches(0.65),
                                   left_w - Inches(0.5),
                                   Inches(4.0))
_add_text_lines(body_l, [
    '·  Feature  : 3 tier — Core(10) / Momentum(14) / Extended(28)',
    '              tier별 독립 튜닝 → "best vs best" 비교',
    '',
    '·  Scaling  : 컬럼 그룹별 (선형 모델만)',
    '              · log_robust (RV/semivariance)',
    '              · robust  (변동률·EPU)',
    '              · standard (거시·모멘텀·spillover)',
    '              · no_scaling (요일·neg_return)',
    '              트리 모델은 raw input',
    '',
    '·  분할       : CSV split 컬럼 (data_splits.txt 정의)',
    '              평시 valid_start = 2012-01-01',
    '              위기 valid = train 끝 ~15%',
    '',
    '·  Target    : RV_target (h=1, 다음 날 RV)',
    '',
    '·  프로토콜 : static (1회 fit → test) + expanding (walk-forward)',
    '              · 금융모형 refit_every = 1',
    '              · ML 모델 refit_every = 5',
], size=11, color=TXT, spacing=2)


# ── 우측: 하이퍼파라미터 튜닝 ──────────
right_left = Inches(7.0)
right_top  = Inches(1.65)
right_w    = Inches(5.6)

label_r = slide.shapes.add_textbox(right_left, right_top, right_w, Inches(0.35))
_set_text(label_r, '2.  Optuna 하이퍼파라미터 튜닝', size=13, bold=True, color=ACCENT)

box_r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                right_left, right_top + Inches(0.45),
                                right_w, Inches(4.3))
box_r.fill.solid()
box_r.fill.fore_color.rgb = ACCENT_LIGHT
box_r.line.color.rgb = ACCENT
box_r.line.width = Pt(1.0)

body_r = slide.shapes.add_textbox(right_left + Inches(0.25),
                                    right_top + Inches(0.65),
                                    right_w - Inches(0.5),
                                    Inches(4.0))
_add_text_lines(body_r, [
    '·  탐색 알고리즘 : Optuna TPE Sampler',
    '                 (seed = 42, 결정론적 탐색)',
    '',
    '·  목적 함수    : valid set의 QLIKE 최소화',
    '',
    '·  n_trials    : 선형 모델 (Ridge / EN / Huber)  =  15',
    '                 트리 모델 (LightGBM / XGBoost)  =  30',
    '',
    '·  탐색 공간 (예시)',
    '   · Ridge      :  alpha ∈ [1e-3, 100]  (log)',
    '   · ElasticNet :  alpha + l1_ratio',
    '   · LightGBM   :  9D (n_est, lr, depth, leaves, ...)',
    '   · XGBoost    :  9D (+ gamma)',
    '',
    '·  튜닝 단위  : (model × regime × country × tier) — 180회',
    '',
    '·  적용 정책 :',
    '   · 1회 튜닝 → best_params 고정',
    '   · static·expanding 모두 동일 best_params 사용',
    '   · expanding refit 시 가중치만 재학습 (hp는 고정)',
], size=11, color=TXT, spacing=2)


# ── 하단: 흐름 요약 ────────────────────
bottom_label = slide.shapes.add_textbox(Inches(0.75), Inches(6.3),
                                          SW - Inches(1.5), Inches(0.35))
_set_text(bottom_label, '전체 흐름', size=12, bold=True, color=ACCENT)

flow = slide.shapes.add_textbox(Inches(0.75), Inches(6.65),
                                  SW - Inches(1.5), Inches(0.5))
_set_text(flow,
    'train (모델 fit)  →  valid (Optuna 튜닝, QLIKE 최소화)  →  best_params 고정  →  '
    'test 평가 (static & expanding)  →  phase별 RMSE / MAE / QLIKE / RMSE_CV',
    size=11, color=TXT_MUTED, italic=True)


# ── 페이지 번호 ────────────────────────
total = len(prs.slides)
page = slide.shapes.add_textbox(SW - Inches(1.2), SH - Inches(0.5),
                                  Inches(1.0), Inches(0.3))
_set_text(page, f'{total} / {total}', size=10, color=TXT_MUTED, align=PP_ALIGN.RIGHT)

# 저장
prs.save(PPTX_PATH)
print(f'saved: {PPTX_PATH}')
print(f'total slides: {total}  (1 슬라이드 추가됨)')

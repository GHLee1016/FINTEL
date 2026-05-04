"""FINTEL 결과 정리 PPT 생성. 깔끔한 plain table 위주, 흰 배경 + 단일 accent."""
from __future__ import annotations
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

PROJECT_ROOT = Path(__file__).resolve().parent.parent

# ── 데이터 로드 ────────────────────────────────────
fin = pd.read_csv(PROJECT_ROOT / 'results/financial_results.csv')
ml_core = pd.read_csv(PROJECT_ROOT / 'results/ml_results_core.csv')
ml_mom  = pd.read_csv(PROJECT_ROOT / 'results/ml_results_momentum.csv')
ml_ext  = pd.read_csv(PROJECT_ROOT / 'results/ml_results_extended.csv')

ml_core['feature_set'] = 'core'
ml_mom['feature_set']  = 'momentum'
ml_ext['feature_set']  = 'extended'
ml_all = pd.concat([ml_core, ml_mom, ml_ext], ignore_index=True)

REGIME_KO = {'normal': '평시', '911': '9·11', 'gfc': 'GFC', 'covid': 'COVID'}
MARKET_FULL = {'US': 'S&P 500 (US)', 'KR': 'KOSPI (KR)', 'JP': 'Nikkei 225 (JP)'}
REGIMES = ['normal', '911', 'gfc', 'covid']
MODEL_ORDER = ['HAR_RV', 'GARCH', 'Ridge', 'ElasticNet', 'Huber', 'LightGBM', 'XGBoost']
TIER_ORDER = ['core', 'momentum', 'extended']


def build_table_models(country: str, protocol: str) -> pd.DataFrame:
    fin_sub = fin[(fin['phase'] == 'Full Test') & (fin['country'] == country)
                  & (fin['protocol'] == protocol)]
    ml_sub = ml_all[(ml_all['phase'] == 'Full Test') & (ml_all['country'] == country)
                    & (ml_all['protocol'] == protocol) & (ml_all['feature_set'] == 'extended')]
    fin_t = fin_sub[['regime', 'model', 'RMSE_CV']]
    ml_t = ml_sub[['regime', 'model', 'RMSE_CV']]
    df = pd.concat([fin_t, ml_t]).pivot(index='regime', columns='model', values='RMSE_CV').round(3)
    df = df.reindex(REGIMES)[MODEL_ORDER]
    return df


def build_table_tier(country: str, protocol: str) -> pd.DataFrame:
    sub = ml_all[(ml_all['phase'] == 'Full Test') & (ml_all['country'] == country)
                 & (ml_all['protocol'] == protocol)]
    agg = sub.groupby(['regime', 'feature_set'])['RMSE_CV'].min().unstack('feature_set')
    agg = agg.reindex(REGIMES)[TIER_ORDER].round(3)
    agg.columns = ['Core', 'Momentum', 'Extended']
    agg['delta_CM'] = (agg['Momentum'] - agg['Core']).round(3)
    agg['delta_ME'] = (agg['Extended'] - agg['Momentum']).round(3)
    agg['pct_CE']   = ((agg['Extended'] / agg['Core'] - 1) * 100).round(1)
    agg.columns = ['Core', 'Momentum', 'Extended', 'D Core->Mom', 'D Mom->Ext', '% Core->Ext']
    return agg


# ── PPT 작성 ─────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ── 색상 (refined Charcoal + Navy accent) ──
ACCENT = RGBColor(0x1E, 0x3A, 0x5F)        # 깊은 네이비 (메인 accent)
ACCENT_LIGHT = RGBColor(0xE8, 0xEE, 0xF4)   # accent의 옅은 톤
HEADER_BG = RGBColor(0x2A, 0x3D, 0x52)      # 표 헤더 (네이비 차콜)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT_BG = RGBColor(0xF7, 0xF8, 0xFA)     # 더 미세한 zebra
TXT = RGBColor(0x1F, 0x2A, 0x3A)            # 거의 검정, 살짝 네이비
TXT_MUTED = RGBColor(0x6B, 0x72, 0x80)      # 캡션·서브
BEST_HL = RGBColor(0xD7, 0xE9, 0xC8)        # 옅은 초록
DIVIDER_GRAY = RGBColor(0xE5, 0xE7, 0xEB)


def _set_text(shape, text, size=14, bold=False, color=TXT, align=PP_ALIGN.LEFT,
              italic=False, font='Calibri'):
    tf = shape.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font


def _add_text_lines(shape, lines, size=14, bold=False, color=TXT, align=PP_ALIGN.LEFT,
                    spacing=4, font='Calibri'):
    tf = shape.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(spacing)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font


def _accent_bar(slide, top, height, color=ACCENT, width_inch=0.08):
    """슬라이드 왼쪽 accent 세로 바."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), top,
                                  Inches(width_inch), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def _hr_line(slide, y, left=Inches(0.5), right_offset=Inches(0.5), color=DIVIDER_GRAY):
    """수평 구분선."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, y,
                                   SW - left - right_offset, Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def _page_number(slide, n, total):
    box = slide.shapes.add_textbox(SW - Inches(1.2), SH - Inches(0.5),
                                    Inches(1.0), Inches(0.3))
    _set_text(box, f'{n} / {total}', size=10, color=TXT_MUTED, align=PP_ALIGN.RIGHT)


def _title_block(slide, title, subtitle=None, accent=True):
    """슬라이드 상단: 타이틀 + 서브타이틀 + 가는 구분선."""
    if accent:
        _accent_bar(slide, Inches(0.45), Inches(0.65))
    tx = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), SW - Inches(1.5), Inches(0.55))
    _set_text(tx, title, size=22, bold=True, color=TXT)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.75), Inches(0.92), SW - Inches(1.5), Inches(0.35))
        _set_text(sub, subtitle, size=12, color=TXT_MUTED, italic=True)
    _hr_line(slide, Inches(1.35))


# ── 표지 슬라이드 ───────────────────────────────────
def add_cover_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    # 큰 accent 막대 — 왼쪽
    _accent_bar(slide, Inches(2.3), Inches(2.8), color=ACCENT, width_inch=0.12)

    # 메인 타이틀
    title = slide.shapes.add_textbox(Inches(0.85), Inches(2.4), SW - Inches(1.7), Inches(1.0))
    _set_text(title, 'FINTEL — 변동성 예측 결과 정리', size=42, bold=True, color=TXT)

    # 부제
    sub = slide.shapes.add_textbox(Inches(0.85), Inches(3.3), SW - Inches(1.7), Inches(0.7))
    _set_text(sub, 'RMSE_CV 비교: ML vs 금융모형 + 데이터셋 추가효과',
              size=18, color=ACCENT)

    # 메타 정보
    meta = slide.shapes.add_textbox(Inches(0.85), Inches(4.5), SW - Inches(1.7), Inches(0.8))
    _add_text_lines(meta,
        ['평시 / 9·11 / GFC / COVID  ×  S&P 500 / KOSPI / Nikkei 225  ×  static / expanding',
         '7 모델 (HAR-RV · GARCH · Ridge · ElasticNet · Huber · LightGBM · XGBoost)'],
        size=13, color=TXT_MUTED, spacing=6)
    return slide


# ── 섹션 헤더 슬라이드 ─────────────────────────────
def add_section_header(prs, label, title, lines):
    slide = prs.slides.add_slide(BLANK)
    # 좌측 큰 accent 영역
    _accent_bar(slide, Inches(2.4), Inches(2.7), width_inch=0.12)

    # 라벨
    lbl = slide.shapes.add_textbox(Inches(0.85), Inches(2.4), SW - Inches(1.7), Inches(0.4))
    _set_text(lbl, label, size=14, bold=True, color=ACCENT)

    # 제목
    ttl = slide.shapes.add_textbox(Inches(0.85), Inches(2.85), SW - Inches(1.7), Inches(0.8))
    _set_text(ttl, title, size=32, bold=True, color=TXT)

    # 본문 (라인들)
    body = slide.shapes.add_textbox(Inches(0.85), Inches(3.85), SW - Inches(1.7), Inches(2.0))
    _add_text_lines(body, lines, size=14, color=TXT_MUTED, spacing=6)
    return slide


# ── RMSE_CV 설명 슬라이드 ──────────────────────────
def add_rmse_cv_explanation(prs):
    slide = prs.slides.add_slide(BLANK)
    _title_block(slide, '왜 RMSE_CV 인가?',
                 '시장·구간 간 공정한 비교를 위한 정규화 지표')

    # 좌: 정의 박스
    box_left = Inches(0.75)
    box_top  = Inches(1.7)
    box_w    = Inches(5.8)
    box_h    = Inches(2.0)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_left, box_top, box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_LIGHT
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.0)

    # 정의 텍스트
    label = slide.shapes.add_textbox(box_left + Inches(0.3), box_top + Inches(0.2),
                                      box_w - Inches(0.6), Inches(0.4))
    _set_text(label, '정의', size=11, bold=True, color=ACCENT)

    formula = slide.shapes.add_textbox(box_left + Inches(0.3), box_top + Inches(0.65),
                                        box_w - Inches(0.6), Inches(0.6))
    _set_text(formula, 'RMSE_CV  =  RMSE  /  mean(y_true)',
              size=22, bold=True, color=TXT, font='Consolas')

    desc = slide.shapes.add_textbox(box_left + Inches(0.3), box_top + Inches(1.35),
                                     box_w - Inches(0.6), Inches(0.6))
    _set_text(desc, 'Coefficient of Variation of RMSE — 무차원 비율',
              size=12, color=TXT_MUTED, italic=True)

    # 우: 왜 쓰는가
    why_top = Inches(1.7)
    why_left = Inches(7.0)
    why_w = Inches(5.8)

    why_label = slide.shapes.add_textbox(why_left, why_top, why_w, Inches(0.4))
    _set_text(why_label, '왜 사용?', size=11, bold=True, color=ACCENT)

    why_body = slide.shapes.add_textbox(why_left, why_top + Inches(0.5), why_w, Inches(2.5))
    _add_text_lines(why_body, [
        '·  시장별 평균 RV 스케일이 다름',
        '   (예: KR 평균 ≈ 0.53,  JP 평균 ≈ 0.83)',
        '',
        '·  같은 RMSE 0.2 라도 KR=38% / JP=24% 의미 다름',
        '',
        '·  RMSE_CV 가 분모로 평균을 나눠 스케일 정규화',
        '   → 시장·구간·모델 간 직접 비교 가능',
    ], size=14, color=TXT, spacing=2)

    # 하단: 수치 감각 표
    table_top = Inches(4.4)
    table_label = slide.shapes.add_textbox(Inches(0.75), table_top, Inches(6), Inches(0.4))
    _set_text(table_label, '수치 감각', size=11, bold=True, color=ACCENT)

    # 작은 표
    rows, cols = 5, 2
    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(0.75), table_top + Inches(0.45),
                                  Inches(8.0), Inches(2.2)).table
    tbl.columns[0].width = Inches(2.0)
    tbl.columns[1].width = Inches(6.0)

    # 헤더
    for j, head in enumerate(['RMSE_CV', '해석']):
        c = tbl.cell(0, j)
        c.text = head
        _style_cell(c, HEADER_BG, HEADER_FG, bold=True, size=12, align=PP_ALIGN.CENTER)

    rows_data = [
        ('0.20', '오차가 평균 RV의 약 20% — 양호한 예측'),
        ('0.30', '오차가 평균의 30% — 보통'),
        ('0.50', '오차가 평균의 50% — 큼'),
        ('1.00', '오차가 평균과 동일 — 모델 무용 수준'),
    ]
    for i, (v, txt) in enumerate(rows_data, start=1):
        c0 = tbl.cell(i, 0); c0.text = v
        _style_cell(c0, ROW_ALT_BG if i % 2 == 0 else None, ACCENT,
                    bold=True, size=14, align=PP_ALIGN.CENTER, font='Consolas')
        c1 = tbl.cell(i, 1); c1.text = txt
        _style_cell(c1, ROW_ALT_BG if i % 2 == 0 else None, TXT,
                    size=13, align=PP_ALIGN.LEFT)

    # 본 PPT 결과 범위 표시
    range_box = slide.shapes.add_textbox(Inches(9.0), table_top + Inches(0.45),
                                          Inches(3.8), Inches(2.2))
    _add_text_lines(range_box, [
        '본 PPT 결과 범위',
        '',
        'ML best :  0.16 ~ 0.34',
        'HAR-RV  :  0.21 ~ 0.50',
        'GARCH   :  0.30 ~ 0.95',
    ], size=13, color=TXT, spacing=4)
    return slide


# ── 데이터 표 슬라이드 ──────────────────────────────
def _style_cell(cell, bg, fg, bold=False, size=11, align=PP_ALIGN.LEFT, font='Calibri'):
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    else:
        cell.fill.background()
    tf = cell.text_frame
    tf.margin_left = Inches(0.07)
    tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = fg
            r.font.name = font


def add_table_slide(prs, title, subtitle, df, caption=None,
                     highlight_min=False, highlight_skip_cols=()):
    slide = prs.slides.add_slide(BLANK)
    _title_block(slide, title, subtitle)

    rows = len(df) + 1
    cols = len(df.columns) + 1
    left = Inches(0.75)
    top = Inches(1.65)
    width = SW - Inches(1.5)
    height = Inches(0.55 + 0.5 * len(df))
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 헤더
    cell = tbl.cell(0, 0)
    cell.text = '구간'
    _style_cell(cell, HEADER_BG, HEADER_FG, bold=True, size=12, align=PP_ALIGN.CENTER)
    for j, col in enumerate(df.columns, start=1):
        cell = tbl.cell(0, j)
        cell.text = str(col)
        _style_cell(cell, HEADER_BG, HEADER_FG, bold=True, size=12, align=PP_ALIGN.CENTER)

    for i, regime in enumerate(df.index, start=1):
        cell = tbl.cell(i, 0)
        cell.text = REGIME_KO.get(regime, regime)
        bg = ROW_ALT_BG if i % 2 == 0 else None
        _style_cell(cell, bg, TXT, bold=True, size=12, align=PP_ALIGN.LEFT)

        if highlight_min:
            row_vals = {}
            for j, col in enumerate(df.columns, start=1):
                if col in highlight_skip_cols:
                    continue
                v = df.iloc[i - 1, j - 1]
                if pd.notna(v):
                    row_vals[j] = v
            best_j = min(row_vals, key=row_vals.get) if row_vals else None
        else:
            best_j = None

        for j, col in enumerate(df.columns, start=1):
            v = df.iloc[i - 1, j - 1]
            cell = tbl.cell(i, j)
            if pd.isna(v):
                cell.text = '—'
            else:
                if col == '% Core->Ext':
                    cell.text = f'{v:+.1f}%' if v != 0 else '0%'
                elif col in ('D Core->Mom', 'D Mom->Ext'):
                    cell.text = f'{v:+.3f}' if v != 0 else '0'
                else:
                    cell.text = f'{v:.3f}'
            row_bg = ROW_ALT_BG if i % 2 == 0 else None
            if best_j == j:
                _style_cell(cell, BEST_HL, ACCENT, bold=True, size=12, align=PP_ALIGN.RIGHT,
                            font='Consolas')
            else:
                _style_cell(cell, row_bg, TXT, bold=False, size=12, align=PP_ALIGN.RIGHT,
                            font='Consolas')

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.75), Inches(6.85),
                                        SW - Inches(1.5), Inches(0.4))
        _set_text(cap, caption, size=10, color=TXT_MUTED, italic=True)
    return slide


# ── 슬라이드 빌드 ──────────────────────────────────
add_cover_slide(prs)

# Section 1
add_section_header(prs, 'SECTION 1', 'ML vs 금융모형 — RMSE_CV',
    ['Full Test phase 기준, ML은 Extended tier 사용',
     '7 모델 비교: HAR-RV · GARCH · Ridge · ElasticNet · Huber · LightGBM · XGBoost',
     '시장 3 (US/KR/JP) × 프로토콜 2 (static/expanding) = 6 표'])

# RMSE_CV 설명 슬라이드 (Section 1 헤더 다음, 첫 표 직전)
add_rmse_cv_explanation(prs)

# 1.1 ~ 1.6
for country in ['US', 'KR', 'JP']:
    for protocol in ['static', 'expanding']:
        df = build_table_models(country, protocol)
        title = f'{MARKET_FULL[country]}  /  {protocol}'
        sub = 'RMSE_CV — Full Test, ML = Extended tier'
        cap = '값이 낮을수록 우수. 각 행 최저값 셀을 옅은 초록으로 강조.'
        add_table_slide(prs, title, sub, df, caption=cap, highlight_min=True)


# ── Section 1 최종 요약: 시장 × 구간 best model ────
def add_best_model_summary(prs):
    """3 시장 × 4 구간 × 2 protocol = 12 행 × (best model + RMSE_CV) 표."""
    slide = prs.slides.add_slide(BLANK)
    _title_block(slide,
        'Section 1 요약 — 시장 × 구간 × Protocol Best Model',
        '각 (시장, 구간, protocol) 셀에서 RMSE_CV 최저 모델 (7 모델 중)')

    # 데이터: 시장×구간×protocol 마다 best 모델
    rows_data = []
    for country in ['US', 'KR', 'JP']:
        for regime in REGIMES:
            row = {'country': country, 'regime': regime}
            for protocol in ['static', 'expanding']:
                df = build_table_models(country, protocol)
                series = df.loc[regime]
                best_m = series.idxmin()
                best_v = series.min()
                row[f'{protocol}_model'] = best_m
                row[f'{protocol}_value'] = best_v
            rows_data.append(row)

    # 표 작성: 6 cols, 13 rows (1 header + 12 data)
    n_rows = len(rows_data) + 1
    n_cols = 6
    left = Inches(0.75); top = Inches(1.55)
    width = SW - Inches(1.5)
    height = Inches(0.55 + 0.4 * len(rows_data))
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # 헤더
    headers = ['시장', '구간', 'static best', 'RMSE_CV', 'expanding best', 'RMSE_CV']
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        _style_cell(c, HEADER_BG, HEADER_FG, bold=True, size=11, align=PP_ALIGN.CENTER)

    # 데이터 (시장별로 첫 행에만 시장명, 나머지는 빈 칸 — merged 효과)
    prev_country = None
    for i, row in enumerate(rows_data, start=1):
        # 시장
        c = tbl.cell(i, 0)
        if row['country'] != prev_country:
            c.text = MARKET_FULL[row['country']]
            _style_cell(c, ROW_ALT_BG if i % 2 == 0 else None, ACCENT,
                         bold=True, size=11, align=PP_ALIGN.LEFT)
            prev_country = row['country']
        else:
            c.text = ''
            _style_cell(c, ROW_ALT_BG if i % 2 == 0 else None, TXT,
                         size=11, align=PP_ALIGN.LEFT)
        # 구간
        c = tbl.cell(i, 1)
        c.text = REGIME_KO[row['regime']]
        _style_cell(c, ROW_ALT_BG if i % 2 == 0 else None, TXT,
                     bold=True, size=11, align=PP_ALIGN.LEFT)
        # static best
        c = tbl.cell(i, 2)
        c.text = row['static_model']
        _style_cell(c, ROW_ALT_BG if i % 2 == 0 else None, ACCENT,
                     bold=True, size=11, align=PP_ALIGN.CENTER)
        c = tbl.cell(i, 3)
        c.text = f'{row["static_value"]:.3f}'
        _style_cell(c, ROW_ALT_BG if i % 2 == 0 else None, TXT,
                     size=11, align=PP_ALIGN.RIGHT, font='Consolas')
        # expanding best
        c = tbl.cell(i, 4)
        c.text = row['expanding_model']
        _style_cell(c, ROW_ALT_BG if i % 2 == 0 else None, ACCENT,
                     bold=True, size=11, align=PP_ALIGN.CENTER)
        c = tbl.cell(i, 5)
        c.text = f'{row["expanding_value"]:.3f}'
        _style_cell(c, ROW_ALT_BG if i % 2 == 0 else None, TXT,
                     size=11, align=PP_ALIGN.RIGHT, font='Consolas')

    # 캡션
    cap = slide.shapes.add_textbox(Inches(0.75), Inches(6.85),
                                    SW - Inches(1.5), Inches(0.4))
    _set_text(cap,
        '7 모델(HAR-RV·GARCH·Ridge·EN·Huber·LGBM·XGB) 중 최저 RMSE_CV 모델. '
        '같은 구간이라도 시장·protocol에 따라 best가 달라짐 — "no free lunch".',
        size=10, color=TXT_MUTED, italic=True)
    return slide


add_best_model_summary(prs)

# Section 2
add_section_header(prs, 'SECTION 2', '데이터셋 추가효과 — Tier 효과',
    ['Best ML RMSE_CV per Tier (Full Test)',
     'Best ML = 5 ML 모델 중 최저 RMSE_CV',
     'Δ는 다음 tier 추가 효과 (음수 = 개선)'])

# 2.1 ~ 2.6
for country in ['US', 'KR', 'JP']:
    for protocol in ['static', 'expanding']:
        df = build_table_tier(country, protocol)
        title = f'{MARKET_FULL[country]}  /  {protocol}'
        sub = 'Tier 효과 — Core(10) → Momentum(14) → Extended(28)'
        cap = ('Best ML = 5 모델(Ridge / EN / Huber / LGBM / XGB) 중 최저 RMSE_CV.  '
               'Δ는 다음 tier 추가 효과 (음수 = 개선).')
        add_table_slide(prs, title, sub, df, caption=cap,
                        highlight_min=True,
                        highlight_skip_cols=('D Core->Mom', 'D Mom->Ext', '% Core->Ext'))

# ── 마지막: Next Steps 슬라이드 ────────────────────
def add_next_steps_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    _title_block(slide,
        'Next Steps — 다음 단계 계획',
        '딥러닝 모형 도입 + 자체 신경망 구조 설계')

    # 좌: 도입 모델
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.65),
                                         Inches(5.8), Inches(0.4))
    _set_text(left_box, '다음 주부터 도입 예정', size=12, bold=True, color=ACCENT)

    models_list = slide.shapes.add_textbox(Inches(0.75), Inches(2.1),
                                            Inches(5.8), Inches(2.5))
    _add_text_lines(models_list, [
        '·  CNN  (Convolutional Neural Network)',
        '·  LSTM  (Long Short-Term Memory)',
        '·  Transformer  (Self-attention 기반)',
        '·  자체 Deep Learning 모델  (FINTEL 설계)',
    ], size=15, color=TXT, spacing=8)

    # 우: 자체 모델 강조 박스
    right_left = Inches(7.0)
    right_top = Inches(1.65)
    right_w = Inches(5.6)
    right_h = Inches(4.5)

    # 강조 박스
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_left, right_top, right_w, right_h)
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_LIGHT
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.2)

    # 박스 내부
    title_in = slide.shapes.add_textbox(right_left + Inches(0.3),
                                         right_top + Inches(0.2),
                                         right_w - Inches(0.6),
                                         Inches(0.4))
    _set_text(title_in, '자체 Deep Learning 모델 — 핵심 설계',
              size=13, bold=True, color=ACCENT)

    body_in = slide.shapes.add_textbox(right_left + Inches(0.3),
                                        right_top + Inches(0.7),
                                        right_w - Inches(0.6),
                                        right_h - Inches(0.9))
    _add_text_lines(body_in, [
        '선행연구를 반영한',
        '그룹구조 신경망 (Group-Structured NN)',
        '',
        '· Feature를 의미별 그룹으로 분리 학습:',
        '    Core (RV/수익률) · Momentum · Macro · Spillover',
        '',
        '· 각 그룹별 sub-network → 후단 통합 layer',
        '',
        '· 기대 효과:',
        '    – feature 그룹별 dynamics 명시적 분리',
        '    – over-parameterization 완화',
        '    – 해석 가능성 향상 (그룹별 contribution)',
    ], size=12, color=TXT, spacing=4)

    # 하단: 본 결과와의 연결
    bottom_label = slide.shapes.add_textbox(Inches(0.75), Inches(6.2),
                                             SW - Inches(1.5), Inches(0.4))
    _set_text(bottom_label, '본 ML 결과와의 연결',
              size=11, bold=True, color=ACCENT)

    bottom_body = slide.shapes.add_textbox(Inches(0.75), Inches(6.55),
                                            SW - Inches(1.5), Inches(0.7))
    _add_text_lines(bottom_body, [
        '· Tier 효과 분석에서 "위기 유형별 feature 그룹 기여도가 다름"을 확인 (예: COVID = 거시, 9·11 = Core only)',
        '· 그룹구조 NN은 이 발견을 아키텍처에 직접 반영 — 단일 fully-connected 대비 더 robust한 학습 기대',
    ], size=10, color=TXT_MUTED, spacing=2)
    return slide


add_next_steps_slide(prs)

# 페이지 번호 (모든 슬라이드 우하단)
total = len(prs.slides)
for n, slide in enumerate(prs.slides, start=1):
    if n == 1:
        continue  # 표지 제외
    _page_number(slide, n, total)

out = PROJECT_ROOT / 'report/fintel_results.pptx'
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f'saved: {out}')
print(f'total slides: {total}')
